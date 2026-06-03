import os
import re
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from docx import Document
from docx.shared import Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn

# ==========================================
# 0. 全局控制器与查重系统状态
# ==========================================
BANNED_WORDS = ["一审", "二圈", "三分段", "四备表达", "四表达", "五成文", "六检查"]

global_seen_pairs = set()
global_last_line_clean = None
global_last_paragraph = None
global_vocab_blacklist = set()

def clean_banned(text):
    if not text: return ""
    text = re.sub(r'[1一]\s*[、\.]?\s*一\s*审|一\s*审', '', text)
    text = re.sub(r'[2二]\s*[、\.]?\s*二\s*圈|二\s*圈', '', text)
    text = re.sub(r'[3三]\s*[、\.]?\s*三\s*分\s*段|三\s*分\s*段', '', text)
    text = re.sub(r'[4四]\s*[、\.]?\s*四\s*备?\s*表\s*达|四\s*备?\s*表\s*达', '', text)
    text = re.sub(r'[5五]\s*[、\.]?\s*五\s*成\s*文|五\s*成\s*文', '', text)
    text = re.sub(r'[6六]\s*[、\.]?\s*六\s*检\s*查|六\s*检\s*查', '', text)
    for w in BANNED_WORDS:
        text = text.replace(w, "")
    return text

# ==========================================
# 1. 基础排版与智能双行查重引擎
# ==========================================
def clean_filename(filename):
    return re.sub(r'[^\w\s\.\-\(\)]', '_', filename)

def set_run_font(run, size_pt, bold=False):
    run.font.name = 'Times New Roman'
    rPr = run._element.get_or_add_rPr()
    rFonts = rPr.get_or_add_rFonts()
    rFonts.set(qn('w:eastAsia'), '宋体')
    rFonts.set(qn('w:ascii'), 'Times New Roman')
    rFonts.set(qn('w:hAnsi'), 'Times New Roman')
    run.font.size = Pt(size_pt)
    run.bold = bold

def add_compact_paragraph(doc, text, bold=False, align=WD_ALIGN_PARAGRAPH.LEFT, color_rgb=None, bypass_blacklist=False, bypass_dup_check=False):
    global global_seen_pairs, global_last_line_clean, global_last_paragraph, global_vocab_blacklist
    
    if not text: return None
        
    if text == "##FORCE_EMPTY##":
        p = doc.add_paragraph()
        p.paragraph_format.line_spacing = 1.5
        global_last_paragraph = None
        return p
        
    if "##INSERT_IMAGE:" in text:
        p = doc.add_paragraph()
        p.paragraph_format.line_spacing = 1.5
        p.alignment = align
        run = p.add_run(text)
        set_run_font(run, 10.5, bold=bold)
        global_last_paragraph = None
        return p

    text_stripped = text.replace('\xa0', ' ').strip()
    if not text_stripped: return None
        
    clean_text = re.sub(r'\s+', '', text_stripped)
    
    if not bypass_blacklist and clean_text in global_vocab_blacklist:
        return None
    
    is_dup = False
    if not bypass_dup_check:
        if global_last_line_clean is not None:
            pair = (global_last_line_clean, clean_text)
            has_chinese = bool(re.search(r'[\u4e00-\u9fa5]', pair[0] + pair[1]))
            
            if has_chinese and pair in global_seen_pairs:
                is_dup = True
                if global_last_paragraph is not None:
                    try:
                        global_last_paragraph._element.getparent().remove(global_last_paragraph._element)
                    except: pass
                    global_last_paragraph = None
            else:
                global_seen_pairs.add(pair)
        global_last_line_clean = clean_text
    else:
        global_last_line_clean = None
        
    if is_dup:
        global_last_paragraph = None
        return None
        
    p = doc.add_paragraph()
    p.paragraph_format.line_spacing = 1.5
    p.paragraph_format.space_before = Pt(0)
    p.paragraph_format.space_after = Pt(0)
    p.alignment = align
    
    run = p.add_run(text_stripped)
    set_run_font(run, 10.5, bold=bold)
    if color_rgb: run.font.color.rgb = color_rgb
        
    global_last_paragraph = p
    return p

# ==========================================
# 2. 空间雷达与颜色识别
# ==========================================
def flatten_shapes(shapes):
    res = []
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP: res.extend(flatten_shapes(shape.shapes))
        else: res.append(shape)
    return res

def get_shape_text(shape):
    if getattr(shape, 'has_text_frame', False): return shape.text
    elif getattr(shape, 'has_table', False): return "".join(cell.text for row in shape.table.rows for cell in row.cells)
    return ""

def get_run_color(run):
    try:
        if run.font and run.font.color and hasattr(run.font.color, 'rgb') and run.font.color.rgb: return str(run.font.color.rgb)
    except: pass
    return None

def is_pink(rgb_str):
    if not rgb_str or rgb_str.lower() == "none": return False
    try:
        r, g, b = int(rgb_str[0:2], 16), int(rgb_str[2:4], 16), int(rgb_str[4:6], 16)
        if r > 150 and g < 180: return True
    except: pass
    return False

def check_target_text_in_slide(slide, target_text):
    for shape in flatten_shapes(slide.shapes):
        if target_text in re.sub(r'\s+', '', get_shape_text(shape)): return True
    return False

def extract_title_bold_only(shape):
    if not getattr(shape, 'has_text_frame', False): return ""
    bold_pieces = []
    for paragraph in shape.text_frame.paragraphs:
        p_txt = ""
        for run in paragraph.runs:
            if run.font.bold: p_txt += run.text
        if p_txt.strip(): bold_pieces.append(clean_banned(p_txt.strip()))
    return " ".join(bold_pieces).strip()

def process_high_freq_paragraphs(shape):
    global global_vocab_blacklist
    local_vocabs = []
    paragraphs = []
    if getattr(shape, 'has_text_frame', False):
        paragraphs = shape.text_frame.paragraphs
    elif getattr(shape, 'has_table', False):
        for row in shape.table.rows:
            for cell in row.cells: paragraphs.extend(cell.text_frame.paragraphs)
            
    if not paragraphs: return local_vocabs
    
    i = 0
    while i < len(paragraphs):
        p = paragraphs[i]
        txt = clean_banned(p.text).strip()
        if not txt: 
            i += 1; continue
            
        is_bold = any(run.font.bold for run in p.runs if run.text.strip())
        is_numbered_eng = bool(re.match(r'^([①-⑳]|\d+[\.、]?)\s*[a-zA-Z]', txt))
        
        if is_numbered_eng and is_bold:
            eng_line = txt
            global_vocab_blacklist.add(re.sub(r'\s+', '', eng_line))
            
            chin_line = ""
            i += 1
            while i < len(paragraphs):
                next_p = paragraphs[i]
                next_txt = clean_banned(next_p.text).strip()
                if not next_txt:
                    i += 1; continue
                if bool(re.match(r'^([①-⑳]|\d+[\.、]?)\s*[a-zA-Z]', next_txt)): break 
                
                chin_line += next_txt + "\n"
                global_vocab_blacklist.add(re.sub(r'\s+', '', next_txt))
                i += 1
                break 
                
            clean_eng = re.sub(r'^([①-⑳]|\d+[\.、]?)\s*', '', eng_line).strip()
            if chin_line: local_vocabs.append(f"{clean_eng}\n{chin_line.strip()}")
            else: local_vocabs.append(clean_eng)
        else:
            i += 1
    return local_vocabs

def stitch_and_bold_vocab(raw_text, pending_number=None):
    raw_text = clean_banned(raw_text)
    raw_text = raw_text.replace('\x0b', '\n').replace('\r', '\n').replace('\xa0', ' ')
    raw_lines = [line.strip() for line in raw_text.split('\n') if line.strip()]
    headers, formatted = [], []
    FILTER_WORDS = ["词汇精讲", "用法搭配", "原文重现", "重点词汇", "文章结构", "参考表达", "必备范文", "知识精讲", "经典题目"]
    
    for line in raw_lines:
        if re.sub(r'\s+', '', line) in FILTER_WORDS: continue
        if pending_number: line, pending_number = f"{pending_number}  {line}", None
        if re.match(r'^([①-⑳]|\d{1,2}[\.、]?)$', line): pending_number = line
        else:
            is_bold = bool(re.match(r'^([①-⑳]|\d{1,2}[\.、]?)\s*', line) and re.search(r'[\u4e00-\u9fa5]', line))
            if is_bold: headers.append(line)
            formatted.append((line, is_bold))
    return headers, formatted, pending_number

# ==========================================
# 4. 英语作文云端调度引擎 (Web 版)
# ==========================================
def process_single_writing_ppt(ppt_path, doc_template, icon_dir, output_path):
    global global_seen_pairs, global_last_line_clean, global_last_paragraph, global_vocab_blacklist
    # 【极度重要】每次运行前必须重置全局状态，防止多次上传产生串扰
    global_seen_pairs = set()
    global_last_line_clean = None
    global_last_paragraph = None
    global_vocab_blacklist = set()
    
    prs = Presentation(ppt_path)
    doc = Document(doc_template)

    slide_sections = {}
    current_sec = "知识精讲"
    title_idx, obj_idx = -1, -1
    high_freq_vocabs = []

    for i, slide in enumerate(prs.slides):
        l_name = slide.slide_layout.name.strip()
        if "标题" in l_name and title_idx == -1: title_idx = i
        if "知识目标" in l_name and obj_idx == -1: obj_idx = i

        found_sec = None
        for shape in flatten_shapes(slide.shapes):
            txt = re.sub(r'\s+', '', get_shape_text(shape))
            if "用法搭配" in txt: found_sec = "词汇精讲"
            elif "文章结构" in txt: found_sec = "文章结构"
            elif "参考表达" in txt: found_sec = "参考表达"
            elif "必备范文" in txt: found_sec = "必备范文"
            elif "经典题目" in txt: found_sec = "经典题目" 
        if found_sec: current_sec = found_sec
        slide_sections[i] = current_sec

        if slide_sections[i] == "参考表达":
            for shape in flatten_shapes(slide.shapes):
                vocabs = process_high_freq_paragraphs(shape)
                if vocabs: high_freq_vocabs.extend(vocabs)

    # --- 模块 1：处理主标题 ---
    if title_idx != -1:
        slide = prs.slides[title_idx]
        title_text = ""
        for shape in slide.shapes:
            t = extract_title_bold_only(shape)
            if t: title_text = t; break
        if title_text:
            merged_title = f"##INSERT_IMAGE:标题块.emf##{title_text}"
            if len(doc.paragraphs) > 0 and not doc.paragraphs[0].text.strip():
                doc.paragraphs[0].text = merged_title
                doc.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
                for run in doc.paragraphs[0].runs: set_run_font(run, 20, bold=True)
            else:
                add_compact_paragraph(doc, merged_title, bold=True, align=WD_ALIGN_PARAGRAPH.CENTER)

    # --- 模块 2：处理课堂目标 ---
    if obj_idx != -1:
        slide = prs.slides[obj_idx]
        add_compact_paragraph(doc, "##INSERT_IMAGE:课堂目标.emf##")
        counter = 1
        shapes = flatten_shapes(slide.shapes)
        shapes.sort(key=lambda s: getattr(s, 'top', 0) or 0)
        for shape in shapes:
            if not getattr(shape, 'has_text_frame', False) or shape == slide.shapes.title: continue 
            for paragraph in shape.text_frame.paragraphs:
                clean_p_txt = clean_banned(paragraph.text)
                if not clean_p_txt.strip(): continue
                p = doc.add_paragraph()
                p.paragraph_format.line_spacing = 1.5
                r_num = p.add_run(f"{counter}. ")
                set_run_font(r_num, 10.5, bold=False)
                for run in paragraph.runs:
                    run_txt = clean_banned(run.text)
                    if not run_txt.strip(): continue
                    is_black = True
                    if run.font and run.font.color and hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                        rgb_str = str(run.font.color.rgb)
                        if rgb_str != "000000" and rgb_str.lower() != "none": is_black = False
                    r_word = p.add_run(run_txt)
                    set_run_font(r_word, 10.5, bold=not is_black)
                counter += 1

    # --- 模块 2.5：课前预习 ---
    if high_freq_vocabs:
        add_compact_paragraph(doc, "##INSERT_IMAGE:课前预习.emf##")
        add_compact_paragraph(doc, "跟读伴学本上的音频，并在讲义横线上填写中文含义。", bold=True)
        for idx, vocab_block in enumerate(high_freq_vocabs):
            lines = vocab_block.split('\n')
            add_compact_paragraph(doc, f"{idx+1}. {lines[0]}", bold=True, bypass_blacklist=True, bypass_dup_check=True)
            for line in lines[1:]:
                add_compact_paragraph(doc, line, bold=True, bypass_blacklist=True, bypass_dup_check=True)
            if idx < len(high_freq_vocabs) - 1:
                add_compact_paragraph(doc, "##FORCE_EMPTY##")

    # --- 模块 3：正文顺序流排版 ---
    last_sec = None
    for i, slide in enumerate(prs.slides):
        if i == title_idx or i == obj_idx: continue
        sec = slide_sections[i]
        
        is_translate = check_target_text_in_slide(slide, "翻译观点句")
        slide_has_underscore = False
        for shape in flatten_shapes(slide.shapes):
            if re.search(r'_{3,}', get_shape_text(shape)): slide_has_underscore = True; break
        should_extract_pink = is_translate or slide_has_underscore
        
        if sec != last_sec:
            if sec == "必备范文" and high_freq_vocabs:
                add_compact_paragraph(doc, "写出下列表达的中文含义，并在写作时借鉴使用。", bold=True, bypass_blacklist=True, bypass_dup_check=True)
                for idx, vocab_block in enumerate(high_freq_vocabs):
                    lines = vocab_block.split('\n')
                    add_compact_paragraph(doc, f"{idx+1}. {lines[0]}", bold=True, bypass_blacklist=True, bypass_dup_check=True)
                    for line in lines[1:]:
                        add_compact_paragraph(doc, line, bold=True, bypass_blacklist=True, bypass_dup_check=True)
            
            if sec in ["词汇精讲", "文章结构", "参考表达", "必备范文", "知识精讲", "经典题目"]:
                add_compact_paragraph(doc, f"##INSERT_IMAGE:{sec}.emf##")
            if sec == "文章结构":
                add_compact_paragraph(doc, "##FORCE_EMPTY##")
            last_sec = sec
            
        elif sec == "词汇精讲" and last_sec == "词汇精讲":
            p_sep = doc.add_paragraph()
            p_sep.paragraph_format.line_spacing = 1.5
            run_sep = p_sep.add_run("﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌")
            set_run_font(run_sep, 10.5)
            p_sep.alignment = WD_ALIGN_PARAGRAPH.CENTER
            add_compact_paragraph(doc, "##INSERT_IMAGE:词汇精讲.emf##")

        if sec == "文章结构": continue

        if sec == "经典题目":
            add_compact_paragraph(doc, "##FORCE_EMPTY##")
            all_classic_lines = []
            shapes = flatten_shapes(slide.shapes)
            shapes.sort(key=lambda x: getattr(x, 'top', 0) or 0)
            for shape in shapes:
                txt_block = clean_banned(get_shape_text(shape))
                if not txt_block.strip(): continue
                lines = txt_block.replace('\x0b', '\n').replace('\r', '\n').split('\n')
                for line in lines:
                    line_clean = line.strip()
                    if not line_clean or re.sub(r'\s+', '', line_clean) == "经典题目": continue
                    all_classic_lines.append(line_clean)
            if all_classic_lines:
                total_lines = len(all_classic_lines)
                for line_idx, line_txt in enumerate(all_classic_lines):
                    if line_idx >= total_lines - 4: add_compact_paragraph(doc, line_txt, bold=False, align=WD_ALIGN_PARAGRAPH.CENTER)
                    else: add_compact_paragraph(doc, line_txt, bold=False, align=WD_ALIGN_PARAGRAPH.LEFT)
            continue
        
        shapes = flatten_shapes(slide.shapes)
        valid_shapes = [s for s in shapes if getattr(s, 'has_text_frame', False) or getattr(s, 'has_table', False)]
        valid_shapes.sort(key=lambda x: (getattr(x, 'top', 0) or 0, getattr(x, 'left', 0) or 0))

        slide_pink_words = []
        shape_data = []

        for shape in valid_shapes:
            if getattr(shape, 'has_table', False):
                table_data = []
                for row in shape.table.rows:
                    row_data = []
                    for cell in row.cells:
                        cell_normal = ""
                        for p in cell.text_frame.paragraphs:
                            p_normal, current_pink = "", ""
                            for run in p.runs:
                                txt = clean_banned(run.text)
                                if not txt: continue
                                is_run_pink = is_pink(get_run_color(run))
                                has_chinese = bool(re.search(r'[\u4e00-\u9fa5]', txt))
                                has_english = bool(re.search(r'[A-Za-z]', txt))
                                
                                if should_extract_pink and is_run_pink and has_english and not has_chinese: current_pink += txt
                                elif should_extract_pink and is_run_pink and not has_english and not has_chinese and current_pink: current_pink += txt
                                else:
                                    if current_pink:
                                        if re.search(r'[A-Za-z]', current_pink): slide_pink_words.append(current_pink.strip())
                                        else: p_normal += current_pink
                                        current_pink = ""
                                    p_normal += txt
                            if current_pink:
                                if re.search(r'[A-Za-z]', current_pink): slide_pink_words.append(current_pink.strip())
                                else: p_normal += current_pink
                            cell_normal += p_normal + "\n"
                        row_data.append(cell_normal.strip())
                    table_data.append(row_data)
                shape_data.append(('table', shape.table, table_data))
            else:
                shape_normal = ""
                for p in shape.text_frame.paragraphs:
                    p_normal, current_pink = "", ""
                    for run in p.runs:
                        txt = clean_banned(run.text)
                        if not txt: continue
                        is_run_pink = is_pink(get_run_color(run))
                        has_chinese = bool(re.search(r'[\u4e00-\u9fa5]', txt))
                        has_english = bool(re.search(r'[A-Za-z]', txt))
                        
                        if should_extract_pink and is_run_pink and has_english and not has_chinese: current_pink += txt
                        elif should_extract_pink and is_run_pink and not has_english and not has_chinese and current_pink: current_pink += txt
                        else:
                            if current_pink:
                                if re.search(r'[A-Za-z]', current_pink): slide_pink_words.append(current_pink.strip())
                                else: p_normal += current_pink
                                current_pink = ""
                            p_normal += txt
                    if current_pink:
                        if re.search(r'[A-Za-z]', current_pink): slide_pink_words.append(current_pink.strip())
                        else: p_normal += current_pink
                    shape_normal += p_normal + "\n"
                shape_data.append(('text', shape_normal, None))

        pending_number = None
        for s_type, obj, tb_data in shape_data:
            if s_type == 'table':
                word_table = doc.add_table(rows=len(obj.rows), cols=len(obj.columns))
                for r_idx, row in enumerate(obj.rows):
                    for c_idx, cell in enumerate(row.cells):
                        is_duplicate = False
                        if r_idx > 0 and cell.text == obj.cell(r_idx-1, c_idx).text and cell.text.strip(): is_duplicate = True
                        if c_idx > 0 and cell.text == obj.cell(r_idx, c_idx-1).text and cell.text.strip(): is_duplicate = True
                        
                        word_cell = word_table.cell(r_idx, c_idx)
                        p = word_cell.paragraphs[0]
                        p.paragraph_format.line_spacing = 1.5
                        
                        if not is_duplicate:
                            cell_text = tb_data[r_idx][c_idx]
                            if not is_translate and slide_pink_words:
                                def repl(match): return slide_pink_words.pop(0) if slide_pink_words else match.group(0)
                                cell_text = re.sub(r'_{3,}', repl, cell_text)
                            
                            if sec in ["词汇精讲", "参考表达"]:
                                headers, formatted, pending_number = stitch_and_bold_vocab(cell_text, pending_number)
                                for i, (text, bold) in enumerate(formatted):
                                    if i > 0:
                                        p = word_cell.add_paragraph()
                                        p.paragraph_format.line_spacing = 1.5
                                    run = p.add_run(text)
                                    set_run_font(run, 10.5, bold=bold)
                            else:
                                raw_lines = cell_text.split('\n')
                                for i, text in enumerate(raw_lines):
                                    text_clean = clean_banned(text).replace('\xa0', ' ').strip()
                                    if not text_clean or re.sub(r'\s+', '', text_clean) in ["知识精讲", "必备范文", "参考表达", "词汇精讲", "文章结构"]: continue
                                    add_compact_paragraph(doc, text_clean)
            else:
                shape_normal = obj
                if not is_translate and slide_pink_words:
                    def repl(match): return slide_pink_words.pop(0) if slide_pink_words else match.group(0)
                    shape_normal = re.sub(r'_{3,}', repl, shape_normal)
                
                if sec in ["词汇精讲", "参考表达"]:
                    _, formatted, pending_number = stitch_and_bold_vocab(shape_normal, pending_number)
                    for line, is_bold in formatted:
                        add_compact_paragraph(doc, line, bold=is_bold)
                else:
                    raw_lines = [l.strip() for l in shape_normal.replace('\x0b','\n').replace('\r','\n').split('\n') if l.strip()]
                    for line in raw_lines:
                        clean_line = re.sub(r'\s+', '', line)
                        if clean_line and clean_line not in ["知识精讲", "必备范文", "参考表达", "词汇精讲", "文章结构"]:
                            add_compact_paragraph(doc, line)
                            
        if pending_number: add_compact_paragraph(doc, pending_number)

        if is_translate and slide_pink_words:
            ans_text = "【答案】" + "；".join(slide_pink_words)
            add_compact_paragraph(doc, ans_text, bold=False, align=WD_ALIGN_PARAGRAPH.LEFT, color_rgb=RGBColor(0, 0, 255))

    # --- 模块 5：全局清理死空行 ---
    while len(doc.paragraphs) > 0 and not doc.paragraphs[0].text.strip() and len(doc.paragraphs[0].runs) == 0:
        p_first = doc.paragraphs[0]._element
        p_first.getparent().remove(p_first)

    consecutive_empty = 0
    p_to_remove = []
    for p in doc.paragraphs:
        if not p.text.strip() and (len(p.runs) == 0 or all(not r.text.strip() for r in p.runs)):
            consecutive_empty += 1
            if consecutive_empty > 2: p_to_remove.append(p)
        else: consecutive_empty = 0
            
    for p in p_to_remove:
        pPr = p._element
        pPr.getparent().remove(pPr)

    doc.save(output_path)