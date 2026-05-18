import os
import io
import re
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from docx import Document
from docx.shared import Emu, Pt, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from PIL import Image

# ==========================================
# 1. 核心配置区
# ==========================================
DOC_TEMPLATE = '模板.docx'
ICON_DIR = 'icons'
MAX_WORD_WIDTH_EMU = 5257800 

GEO_ICON_MAP = {
    "知识目标": "知识目标.emf",
    "知识清单": "知识清单.emf",
    "巩固练习": "不插入",
    "经典例题": "经典例题.emf",
    "笔记总结": "笔记总结.emf"
}

# ==========================================
# 2. 增强型排版引擎
# ==========================================
def clean_filename(filename):
    return re.sub(r'[^\w\s\.\-\(\)]', '_', filename)

def set_run_font(run, size_pt, bold=False):
    run.font.name = 'Times New Roman'
    rPr = run._element.get_or_add_rPr()
    rFonts = rPr.get_or_add_rFonts()
    rFonts.set(qn('w:eastAsia'), '宋体')
    rFonts.set(qn('w:ascii'), 'Times New Roman')
    rFonts.set(qn('w:hAnsi'), 'Times New Roman')
    run.font.size = Pt(size_pt)
    run.bold = bold

def add_geo_paragraph(doc, text, text_type='body', is_first=False):
    text = str(text).replace('\xa0', ' ').strip()
    if not text: return None
    text = re.sub(r'_{3,}', '______', text)

    if is_first and len(doc.paragraphs) > 0:
        p = doc.paragraphs[0]
        p.text = "" 
    else:
        p = doc.add_paragraph()

    p.paragraph_format.line_spacing = 1.5 
    p.paragraph_format.space_before = Pt(0)
    p.paragraph_format.space_after = Pt(0)
    
    run = p.add_run(text)

    if text_type == 'title_page':
        set_run_font(run, 22, bold=True)       
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER 
    elif text_type == 'h1':  
        set_run_font(run, 12, bold=True)       
        p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    elif text_type == 'h2':  
        set_run_font(run, 10.5, bold=False)    
        p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    elif text_type == 'caption':
        set_run_font(run, 9, bold=False)       
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.line_spacing = 1.0 
    else:  
        set_run_font(run, 10.5, bold=False)    
        p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
        
    return p

def set_table_borders(table):
    tbl = table._element
    tblPr = tbl.xpath('w:tblPr')
    if not tblPr:
        tblPr = OxmlElement('w:tblPr')
        tbl.insert(0, tblPr)
    else:
        tblPr = tblPr[0]
        
    tblBorders = tblPr.xpath('w:tblBorders')
    if not tblBorders:
        tblBorders = OxmlElement('w:tblBorders')
        tblPr.append(tblBorders)
    else:
        tblBorders = tblBorders[0]
        tblBorders.clear()
        
    for b_name in ['top', 'left', 'bottom', 'right', 'insideH', 'insideV']:
        b = OxmlElement(f'w:{b_name}')
        b.set(qn('w:val'), 'single')
        b.set(qn('w:sz'), '4') 
        b.set(qn('w:space'), '0')
        b.set(qn('w:color'), 'auto')
        tblBorders.append(b)

def flatten_shapes(shapes):
    res = []
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            res.extend(flatten_shapes(shape.shapes))
        else:
            res.append(shape)
    return res

def get_true_cropped_image(shape):
    blob = shape.image.blob
    ext = shape.image.ext.lower()
    if ext in ['emf', 'wmf', 'gif']: return blob
    try:
        cl, cr, ct, cb = getattr(shape, 'crop_left', 0), getattr(shape, 'crop_right', 0), getattr(shape, 'crop_top', 0), getattr(shape, 'crop_bottom', 0)
        if sum([cl, cr, ct, cb]) > 0.001:
            img = Image.open(io.BytesIO(blob))
            w, h = img.size
            cropped = img.crop((cl*w, ct*h, w-(cr*w), h-(cb*h)))
            out_io = io.BytesIO()
            fmt = img.format if img.format else 'PNG'
            if fmt == 'JPEG' and cropped.mode in ('RGBA', 'P'): cropped = cropped.convert('RGB')
            cropped.save(out_io, format=fmt)
            return out_io.getvalue()
    except: pass
    return blob

def identify_real_section(slide):
    try:
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title_text = slide.shapes.title.text.strip()
            for key in GEO_ICON_MAP.keys():
                if key in title_text: return key
    except: pass
    layout_name = slide.slide_layout.name.strip()
    for key in GEO_ICON_MAP.keys():
        if key in layout_name: return key
    return layout_name

def is_red_font(run):
    if hasattr(run, 'font') and run.font.color and hasattr(run.font.color, 'rgb') and run.font.color.rgb:
        rgb_str = str(run.font.color.rgb)
        if len(rgb_str) == 6:
            try:
                r = int(rgb_str[0:2], 16)
                g = int(rgb_str[2:4], 16)
                b = int(rgb_str[4:6], 16)
                if r > 120 and r > g + 30 and r > b + 30: return True
            except: pass
    return False

def extract_text_and_answers(item, is_exercise_slide):
    if hasattr(item, 'has_text_frame') and not item.has_text_frame: return "", []
    if not hasattr(item, 'text_frame') or not item.text_frame: return "", []
    
    normal_texts, red_answers = [], []
    for paragraph in item.text_frame.paragraphs:
        p_normal, current_red_ans = "", ""
        for run in paragraph.runs:
            if is_red_font(run):
                current_red_ans += run.text
                if not is_exercise_slide: p_normal += run.text 
            else:
                if current_red_ans.strip():
                    red_answers.append(current_red_ans.strip())
                    current_red_ans = ""
                p_normal += run.text
        if current_red_ans.strip(): red_answers.append(current_red_ans.strip())
        if p_normal.strip(): normal_texts.append(p_normal.strip())
    return "\n".join(normal_texts), red_answers

# ==========================================
# 3. 核心提取逻辑 (已加入 output_path 参数)
# ==========================================
def process_single_geo_ppt(ppt_path, doc_template, icon_dir, output_path=None):
    raw_filename = os.path.splitext(os.path.basename(ppt_path))[0]
    safe_filename = clean_filename(raw_filename)
    
    if output_path is None:
        output_path = f"{raw_filename}地理讲义.docx"
    
    prs = Presentation(ppt_path)
    doc = Document(doc_template)
    
    last_layout_name = None
    recent_headings = [] 
    last_h2_prefix_clean = None 
    last_h2_paragraph = None 
    
    total_slides = len(prs.slides)
    content_started = False 

    print(f"\n🚀 正在处理: [{raw_filename}] (共 {total_slides} 页)")

    for slide_number, slide in enumerate(prs.slides, 1):
        if slide_number == total_slides: continue

        if slide_number == 1:
            main_title = ""
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                main_title = slide.shapes.title.text.strip()
            if not main_title: 
                text_shapes = [s for s in slide.shapes if s.has_text_frame and s.text.strip()]
                if text_shapes:
                    text_shapes.sort(key=lambda s: s.width * s.height, reverse=True)
                    main_title = text_shapes[0].text.strip()
            
            if main_title:
                add_geo_paragraph(doc, f"##INSERT_IMAGE:标题块.emf##{main_title}", 'title_page', is_first=True)
                content_started = True
            continue

        real_section_name = identify_real_section(slide)

        if real_section_name != last_layout_name:
            icon_file = GEO_ICON_MAP.get(real_section_name, "不插入")
            if icon_file != "不插入":
                p = add_geo_paragraph(doc, f"##INSERT_IMAGE:{icon_file}##", 'body', is_first=not content_started)
                if p: 
                    p.alignment = WD_ALIGN_PARAGRAPH.LEFT
                    content_started = True

        if real_section_name == "笔记总结":
            last_layout_name = real_section_name
            continue 

        all_elements = flatten_shapes(slide.shapes)
        all_elements.sort(key=lambda s: (s.top if hasattr(s, 'top') and s.top is not None else 0))
        processed_shape_ids = set()
        
        slide_red_answers = []
        is_exercise_slide = (real_section_name in ["经典例题", "巩固练习"])

        for i, shape in enumerate(all_elements):
            if id(shape) in processed_shape_ids: continue

            if getattr(shape, 'has_table', False):
                ppt_table = shape.table
                word_table = doc.add_table(rows=len(ppt_table.rows), cols=len(ppt_table.columns))
                
                set_table_borders(word_table)
                
                for r_idx, row in enumerate(ppt_table.rows):
                    for c_idx, cell in enumerate(row.cells):
                        
                        is_duplicate = False
                        if r_idx > 0 and cell.text == ppt_table.cell(r_idx-1, c_idx).text and cell.text.strip():
                            is_duplicate = True
                        if c_idx > 0 and cell.text == ppt_table.cell(r_idx, c_idx-1).text and cell.text.strip():
                            is_duplicate = True
                            
                        word_cell = word_table.cell(r_idx, c_idx)
                        p = word_cell.paragraphs[0]
                        p.text = ""
                        p.paragraph_format.line_spacing = 1.5
                        p.paragraph_format.space_before = Pt(0)
                        p.paragraph_format.space_after = Pt(0)
                        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        
                        if not is_duplicate:
                            normal_text, red_answers = extract_text_and_answers(cell, is_exercise_slide)
                            if is_exercise_slide and red_answers: 
                                slide_red_answers.extend(red_answers)
                            
                            text = normal_text.strip()
                            if text:
                                run = p.add_run(text)
                                set_run_font(run, 10.5, bold=False)
                                
                processed_shape_ids.add(id(shape))
                content_started = True
                continue

            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_blob = get_true_cropped_image(shape)
                sw, sh = shape.width, shape.height
                
                caption_text = None
                if i + 1 < len(all_elements):
                    next_s = all_elements[i+1]
                    if hasattr(next_s, 'has_text_frame') and next_s.has_text_frame and len(next_s.text.strip()) < 30 and "\n" not in next_s.text:
                        if next_s.top > (shape.top + shape.height * 0.5) and \
                           not (next_s.left + next_s.width < shape.left or next_s.left > shape.left + shape.width):
                            _, test_red = extract_text_and_answers(next_s, True)
                            if not test_red:
                                caption_text = next_s.text.strip()
                                processed_shape_ids.add(id(next_s))

                if sw > MAX_WORD_WIDTH_EMU:
                    sh = int(sh * (MAX_WORD_WIDTH_EMU / sw))
                    sw = MAX_WORD_WIDTH_EMU
                
                p_img = doc.add_paragraph()
                p_img.alignment = WD_ALIGN_PARAGRAPH.CENTER
                p_img.paragraph_format.space_before = Pt(0)
                p_img.paragraph_format.space_after = Pt(0)
                r_img = p_img.add_run()
                
                if shape.image.ext.lower() in ['emf', 'wmf']:
                    img_name = f"{safe_filename}_s{slide_number}_{shape.shape_id}.{shape.image.ext.lower()}"
                    with open(os.path.join(icon_dir, img_name), 'wb') as f: f.write(img_blob)
                    r_img.add_text(f"##INSERT_IMAGE:{img_name}|{sw/12700:.1f}|{sh/12700:.1f}##")
                else:
                    r_img.add_picture(io.BytesIO(img_blob), width=Emu(sw), height=Emu(sh))
                
                if caption_text: add_geo_paragraph(doc, caption_text, 'caption')
                processed_shape_ids.add(id(shape))
                content_started = True

            elif hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                normal_text, red_answers = extract_text_and_answers(shape, is_exercise_slide)
                if is_exercise_slide and red_answers: slide_red_answers.extend(red_answers)
                text = normal_text.strip()
                if not text: processed_shape_ids.add(id(shape)); continue 
                
                pure_text = re.sub(r'\s+', '', text)
                merge_match = re.match(r'^([一二三四五六七八九十]+、[^\s①-⑩]*|\d+[\.、．][^\s①-⑩]+)\s*([①②③④⑤⑥⑦⑧⑨⑩].*)', text.replace('\n', ''))
                
                if merge_match:
                    prefix_clean = re.sub(r'\s+', '', merge_match.group(1).strip())
                    if prefix_clean == last_h2_prefix_clean and last_h2_paragraph is not None:
                        run = last_h2_paragraph.add_run("\n" + merge_match.group(2).strip())
                        set_run_font(run, 10.5, bold=False)
                        processed_shape_ids.add(id(shape)); continue
                    else:
                        last_h2_prefix_clean = prefix_clean
                        p_new = add_geo_paragraph(doc, text, 'h2', is_first=not content_started)
                        last_h2_paragraph, content_started = p_new, True
                else:
                    is_h1 = (shape == slide.shapes.title) or bool(re.match(r'^([一二三四五六七八九十]+、)', text))
                    is_h2 = bool(re.match(r'^(\d+[\.、．])', text))
                    t_type = 'h1' if is_h1 else ('h2' if is_h2 else 'body')
                    if t_type in ['h1', 'h2']:
                        if pure_text in recent_headings: processed_shape_ids.add(id(shape)); continue 
                        recent_headings.append(pure_text); (recent_headings.pop(0) if len(recent_headings) > 5 else None)
                    p_normal = add_geo_paragraph(doc, text, t_type, is_first=not content_started)
                    if p_normal: 
                        content_started = True
                        if t_type in ['h1', 'h2']:
                            pre = re.match(r'^([一二三四五六七八九十]+、[^\s]+|\d+[\.、．][^\s]+)', text)
                            if pre: last_h2_prefix_clean, last_h2_paragraph = re.sub(r'\s+', '', pre.group(1)), p_normal
                processed_shape_ids.add(id(shape))

        if is_exercise_slide and slide_red_answers:
            p_ans = add_geo_paragraph(doc, f"【答案】{'；'.join(slide_red_answers)}", 'body', is_first=not content_started)
            if p_ans: 
                for r in p_ans.runs: r.bold = True
                content_started = True
        last_layout_name = real_section_name

    doc.save(output_path)
    print(f"✅ 地理讲义处理完成: {output_path}\n")

if __name__ == '__main__':
    if not os.path.exists(ICON_DIR): os.makedirs(ICON_DIR)
    files = [f for f in os.listdir('.') if f.lower().endswith('.pptx') and not f.startswith('~$')]
    for f in files: process_single_geo_ppt(f, DOC_TEMPLATE, ICON_DIR)
