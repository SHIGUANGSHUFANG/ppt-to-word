import os
import io
import re
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from docx import Document
from docx.shared import Emu, Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_TAB_ALIGNMENT
from docx.oxml.ns import qn
from PIL import Image

# ==========================================
# 1. 核心配置区
# ==========================================
DOC_TEMPLATE = '模板.docx'
ICON_DIR = 'icons'
MAX_WORD_WIDTH_EMU = 5257800 

ENGLISH_ICON_MAP = {
    "标题": "标题块.emf",
    "知识目标": "课堂目标.emf",
    "知识精讲": "知识精讲.emf",
    "词汇精讲": "词汇精讲.emf"
}

QUESTION_TAGS = ["细节理解", "句子还原", "词义猜测", "主旨大意"]

# ==========================================
# 2. 文本与排版净化引擎
# ==========================================
def clean_xml_text(text):
    if not text: return ""
    return re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)

def normalize_for_dedup(text):
    if not text: return ""
    return re.sub(r'[^\w\u4e00-\u9fa5]', '', text.lower())

def clean_filename(filename):
    return re.sub(r'[^\w\s\.\-\(\)]', '_', filename)

def set_run_font(run, size_pt, bold=False):
    run.font.name = 'Times New Roman'
    rPr = run._element.get_or_add_rPr()
    rFonts = rPr.get_or_add_rFonts()
    rFonts.set(qn('w:eastAsia'), '宋体')
    rFonts.set(qn('w:ascii'), 'Times New Roman')
    rFonts.set(qn('w:hAnsi'), 'Times New Roman')
    run.font.size = Pt(size_pt)
    run.bold = bold

def add_compact_paragraph(doc_or_cell, text, text_type='body', is_first=False):
    text = clean_xml_text(str(text).replace('\xa0', ' ').strip())
    if not text: return None

    if is_first and hasattr(doc_or_cell, 'paragraphs') and len(doc_or_cell.paragraphs) > 0:
        p = doc_or_cell.paragraphs[0]
        p.text = ""
    else:
        p = doc_or_cell.add_paragraph()

    p.paragraph_format.line_spacing = 1.5 
    p.paragraph_format.space_before = Pt(0)
    p.paragraph_format.space_after = Pt(0)
    
    run = p.add_run(text)

    if text_type == 'title_page':
        set_run_font(run, 20, bold=True)       
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER 
    elif text_type == 'h1':  
        set_run_font(run, 10.5, bold=True)     
        p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    else:  
        set_run_font(run, 10.5, bold=False)    
        p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
        
    return p

def add_runs_with_blue_tags(paragraph, text):
    pattern = r'(细节理解|句子还原|词义猜测|主旨大意)'
    parts = re.split(pattern, text)
    for part in parts:
        if part in QUESTION_TAGS:
            r = paragraph.add_run(part)
            set_run_font(r, 10.5, bold=False)
            r.font.color.rgb = RGBColor(0, 0, 255) 
        elif part:
            r = paragraph.add_run(part)
            set_run_font(r, 10.5, bold=False)

# ==========================================
# 3. 几何与底层XML解析雷达
# ==========================================
def flatten_shapes(shapes):
    res = []
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            res.extend(flatten_shapes(shape.shapes))
        else:
            res.append(shape)
    return res

def extract_text_safe(shape):
    if getattr(shape, 'has_text_frame', False) and shape.text.strip():
        return shape.text.strip()
    return ""

def extract_title_bold_only(shape):
    if not getattr(shape, 'has_text_frame', False): return ""
    bold_pieces = []
    for paragraph in shape.text_frame.paragraphs:
        p_txt = ""
        for run in paragraph.runs:
            if run.font.bold: p_txt += run.text
        if p_txt.strip(): bold_pieces.append(p_txt.strip())
    return clean_xml_text(" ".join(bold_pieces))

def get_slide_sections(prs):
    """【黑魔法】强行破解 PPT 底层 XML，提取每个幻灯片所属的“节”名称"""
    slide_sections = {}
    try:
        # 寻找所有的节节点
        sections = prs.element.xpath('.//*[local-name()="section"]')
        slide_id_to_section = {}
        for section in sections:
            name = section.get('name')
            # 提取该节下所有的幻灯片ID
            for sld in section.xpath('.//*[local-name()="sldId"]'):
                val = sld.get('id')
                if val:
                    slide_id_to_section[int(val)] = name
                    
        # 将幻灯片编号与节名称挂钩
        for idx, slide in enumerate(prs.slides, 1):
            slide_sections[idx] = slide_id_to_section.get(slide.slide_id, "")
    except Exception as e:
        print(f"⚠️ 解析幻灯片节时出现异常: {e}")
    return slide_sections

# ==========================================
# 4. 核心调度大脑
# ==========================================
def process_single_english_ppt(ppt_path, doc_template, icon_dir, output_path=None):
    raw_filename = os.path.splitext(os.path.basename(ppt_path))[0]
    safe_filename = clean_filename(raw_filename)
    if output_path is None:
        output_path = f"{raw_filename}英语讲义.docx"
    
    prs = Presentation(ppt_path)
    doc = Document(doc_template)
    
    # 获取每一页幻灯片所属的“节(Section)”名称
    slide_sections = get_slide_sections(prs)
    
    last_active_name = None
    content_started = False 

    # --------------------------------------------------
    # 【第一阶段】（Pass 1）：预建去重池与提纯词汇解析
    # --------------------------------------------------
    knowledge_paragraphs_pool = set()
    
    # 1.1 先遣队 1：扫荡“知识精讲”，构建全局去重记忆库
    for slide_number, slide in enumerate(prs.slides, 1):
        sec_name = slide_sections.get(slide_number, "")
        layout_name = slide.slide_layout.name.strip()
        # 【双擎驱动】：优先使用幻灯片节名称，如果没有或不匹配，则降级使用版式名称
        active_name = sec_name if sec_name in ENGLISH_ICON_MAP else layout_name
        
        if active_name == "知识精讲":
            all_elements = flatten_shapes(slide.shapes)
            all_elements.sort(key=lambda s: getattr(s, 'top', 0) or 0)
            
            has_28pt = False
            text_28pt_contents = []
            for shape in all_elements:
                if getattr(shape, 'has_text_frame', False):
                    shape_28_text = ""
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font and run.font.size and round(run.font.size.pt) == 28:
                                has_28pt = True
                                shape_28_text += run.text
                        shape_28_text += "\n"
                    clean_shape_text = clean_xml_text(shape_28_text.strip())
                    if clean_shape_text:
                        text_28pt_contents.append(clean_shape_text)
                        
            if has_28pt:
                for txt in text_28pt_contents:
                    knowledge_paragraphs_pool.add(normalize_for_dedup(txt))
            else:
                valid_text_shapes = [s for s in all_elements if getattr(s, 'has_text_frame', False)]
                for shape in valid_text_shapes:
                    text_content = clean_xml_text(extract_text_safe(shape))
                    if text_content:
                        knowledge_paragraphs_pool.add(normalize_for_dedup(text_content))

    # 1.2 先遣队 2：在“词汇精讲”中严格执行“先过滤，再提取”
    vocab_database = []
    for slide_number, slide in enumerate(prs.slides, 1):
        sec_name = slide_sections.get(slide_number, "")
        layout_name = slide.slide_layout.name.strip()
        active_name = sec_name if sec_name in ENGLISH_ICON_MAP else layout_name
        
        if active_name == "词汇精讲":
            slide_lines = []
            v_elements = flatten_shapes(slide.shapes)
            v_elements.sort(key=lambda s: getattr(s, 'top', 0) or 0)
            
            for shape in v_elements:
                if getattr(shape, 'has_text_frame', False) and shape.text.strip():
                    for paragraph in shape.text_frame.paragraphs:
                        txt = clean_xml_text(paragraph.text.strip())
                        if not txt: continue
                        
                        clean_para_text = normalize_for_dedup(txt)
                        sentences_count = len(re.findall(r'[.!?。！？]', txt))
                        if sentences_count >= 2 or len(clean_para_text) > 60:
                            is_duplicate = False
                            for pool_text in knowledge_paragraphs_pool:
                                if clean_para_text in pool_text or pool_text in clean_para_text:
                                    is_duplicate = True
                                    break
                            if is_duplicate:
                                continue 
                                
                        slide_lines.append(txt)
            
            v_idx = 0
            while v_idx < len(slide_lines):
                line = slide_lines[v_idx]
                if re.match(r'^\s*\d+\s*[a-zA-Z]', line):
                    vocab_text = line
                    translation_text = ""
                    
                    search_idx = v_idx + 1
                    while search_idx < len(slide_lines):
                        next_line = slide_lines[search_idx]
                        if re.match(r'^\s*\d+\s*[a-zA-Z]', next_line):
                            break  
                        
                        if re.search(r'[\u4e00-\u9fa5]', next_line):
                            translation_text = next_line
                            break
                        search_idx += 1
                    
                    vocab_database.append({
                        "vocab": vocab_text,
                        "translation": translation_text
                    })
                v_idx += 1

    total_slides = len(prs.slides)
    print(f"\n🚀 英语讲义引擎启动: [{raw_filename}] (共 {total_slides} 页)")

    # --------------------------------------------------
    # 【第二阶段】（Pass 2）：标准版面渲染与排版组装
    # --------------------------------------------------
    for slide_number, slide in enumerate(prs.slides, 1):
        # 实时判定当前页面的主导名称（幻灯片节优先，版式名称托底）
        sec_name = slide_sections.get(slide_number, "")
        layout_name = slide.slide_layout.name.strip()
        active_name = sec_name if sec_name in ENGLISH_ICON_MAP else layout_name
        
        if active_name != "标题" and active_name not in ENGLISH_ICON_MAP:
            continue
            
        print(f"   正在处理第 {slide_number} 页: [{active_name}]")

        if active_name == "标题":
            title_text = ""
            for shape in slide.shapes:
                extracted = extract_title_bold_only(shape)
                if extracted: title_text = extracted; break
            
            if title_text:
                add_compact_paragraph(doc, f"##INSERT_IMAGE:标题块.emf##{title_text}", 'title_page', is_first=True)
                content_started = True
            continue

        if active_name != last_active_name:
            icon_file = ENGLISH_ICON_MAP.get(active_name, "不插入")
            if icon_file != "不插入":
                add_compact_paragraph(doc, f"##INSERT_IMAGE:{icon_file}##", 'h1', is_first=not content_started)
                content_started = True

        all_elements = flatten_shapes(slide.shapes)
        all_elements.sort(key=lambda s: getattr(s, 'top', 0) or 0)
        
        # --- 知识目标板块 ---
        if active_name == "知识目标":
            line_counter = 1
            for shape in all_elements:
                if getattr(shape, 'has_text_frame', False) and shape.text.strip():
                    for paragraph in shape.text_frame.paragraphs:
                        if not paragraph.text.strip(): continue
                        
                        p_final = doc.add_paragraph()
                        p_final.paragraph_format.line_spacing = 1.5
                        p_final.paragraph_format.space_before = Pt(0)
                        p_final.paragraph_format.space_after = Pt(0)
                        
                        run_num = p_final.add_run(f"{line_counter}. ")
                        set_run_font(run_num, 10.5, bold=False)
                        
                        for run in paragraph.runs:
                            r_txt = clean_xml_text(run.text)
                            if not r_txt: continue
                            
                            is_black = True
                            if run.font and run.font.color and hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                                rgb_str = str(run.font.color.rgb)
                                if rgb_str != "000000" and rgb_str.lower() != "none":
                                    is_black = False 
                            
                            r_word = p_final.add_run(r_txt)
                            set_run_font(r_word, 10.5, bold=(not is_black))
                        
                        line_counter += 1
            
            # --- 课前预习空降逻辑 ---
            if vocab_database:
                add_compact_paragraph(doc, "##INSERT_IMAGE:课前预习.emf##", 'h1')
                
                p_guide = doc.add_paragraph()
                p_guide.paragraph_format.line_spacing = 1.5
                p_guide.paragraph_format.space_before = Pt(0)
                p_guide.paragraph_format.space_after = Pt(0)
                r_guide = p_guide.add_run("跟读伴学本上的音频，并在讲义横线上填写中文含义。")
                set_run_font(r_guide, 10.5, bold=True)
                
                for item in vocab_database:
                    p_v = doc.add_paragraph()
                    p_v.paragraph_format.line_spacing = 1.5
                    p_v.paragraph_format.space_before = Pt(0)
                    p_v.paragraph_format.space_after = Pt(0)
                    
                    r_v_text = p_v.add_run(item["vocab"])
                    set_run_font(r_v_text, 10.5, bold=True) 
                    
                    if item["translation"]:
                        p_v.add_run("    ") 
                        r_t_text = p_v.add_run(item["translation"])
                        set_run_font(r_t_text, 10.5, bold=False)
                        r_t_text.font.color.rgb = RGBColor(0, 0, 255) 
                    
                    p_ex = doc.add_paragraph()
                    p_ex.paragraph_format.line_spacing = 1.5
                    p_ex.paragraph_format.space_before = Pt(0)
                    p_ex.paragraph_format.space_after = Pt(0)
                    p_ex.add_run("") 
            
            last_active_name = active_name
            continue

        # --- 知识精讲板块 ---
        if active_name == "知识精讲":
            has_28pt = False
            text_28pt_contents = []
            
            for shape in all_elements:
                if getattr(shape, 'has_text_frame', False):
                    shape_28_text = ""
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font and run.font.size and round(run.font.size.pt) == 28:
                                has_28pt = True
                                shape_28_text += run.text
                        shape_28_text += "\n"
                    
                    clean_shape_text = clean_xml_text(shape_28_text.strip())
                    if clean_shape_text:
                        text_28pt_contents.append(clean_shape_text)
            
            if has_28pt:
                for txt in text_28pt_contents:
                    p_normal = doc.add_paragraph()
                    p_normal.paragraph_format.line_spacing = 1.5
                    p_normal.paragraph_format.space_before = Pt(0)
                    p_normal.paragraph_format.space_after = Pt(0)
                    
                    lines = txt.split('\n')
                    for l_idx, l_text in enumerate(lines):
                        add_runs_with_blue_tags(p_normal, l_text)
                        if l_idx < len(lines) - 1:
                            p_normal.add_run('\n')
                            
                last_active_name = active_name
                continue

            valid_text_shapes = [s for s in all_elements if getattr(s, 'has_text_frame', False)]
            tag_map = {} 
            skip_indices = set()
            
            for i, shape in enumerate(valid_text_shapes):
                text = clean_xml_text(extract_text_safe(shape)).strip()
                for tag in QUESTION_TAGS:
                    if tag in text and len(text) < 15:
                        closest_idx, min_dist = -1, float('inf')
                        for j, other_shape in enumerate(valid_text_shapes):
                            if i == j: continue
                            
                            other_text = clean_xml_text(extract_text_safe(other_shape)).strip()
                            is_other_tag = any(t in other_text and len(other_text) < 15 for t in QUESTION_TAGS)
                            if is_other_tag: continue
                            
                            dy = abs(getattr(shape, 'top', 0) - getattr(other_shape, 'top', 0))
                            
                            if re.match(r'^\s*\d+[\.、．]', other_text):
                                dy -= 800000 
                            
                            if dy < 3000000: 
                                if dy < min_dist:
                                    min_dist, closest_idx = dy, j
                        
                        if closest_idx != -1:
                            if closest_idx not in tag_map:
                                tag_map[closest_idx] = []
                            tag_map[closest_idx].append(f"【{tag}题】")
                            skip_indices.add(i)
                        break

            for idx, shape in enumerate(valid_text_shapes):
                if idx in skip_indices: continue
                
                text_content = clean_xml_text(extract_text_safe(shape))
                if not text_content: continue
                
                p_normal = doc.add_paragraph()
                p_normal.paragraph_format.line_spacing = 1.5
                p_normal.paragraph_format.space_before = Pt(0)
                p_normal.paragraph_format.space_after = Pt(0)
                
                tab_stops = p_normal.paragraph_format.tab_stops
                tab_stops.add_tab_stop(Cm(15.5), WD_TAB_ALIGNMENT.RIGHT)
                
                if idx in tag_map:
                    tags = tag_map[idx] 
                    lines = text_content.split('\n')
                    
                    tag_idx = 0
                    for line_num, line in enumerate(lines):
                        if tag_idx >= len(tags): break
                        if len(tags) == 1 and line_num == 0:
                            lines[line_num] = line + "\t" + tags[tag_idx]
                            tag_idx += 1
                        elif len(tags) > 1 and re.match(r'^\s*\d+[\.、．]', line):
                            lines[line_num] = line + "\t" + tags[tag_idx]
                            tag_idx += 1
                    
                    while tag_idx < len(tags):
                        lines[0] = lines[0] + "\t" + tags[tag_idx]
                        tag_idx += 1
                    
                    text_content = '\n'.join(lines)
                
                lines = text_content.split('\n')
                for l_idx, l_text in enumerate(lines):
                    add_runs_with_blue_tags(p_normal, l_text)
                    if l_idx < len(lines) - 1:
                        p_normal.add_run('\n')
            
            last_active_name = active_name
            continue

        # --- 词汇精讲板块 ---
        if active_name == "词汇精讲":
            slide_lines = []
            for shape in all_elements:
                if getattr(shape, 'has_text_frame', False) and shape.text.strip():
                    for paragraph in shape.text_frame.paragraphs:
                        txt = clean_xml_text(paragraph.text.strip())
                        if not txt: continue
                        
                        clean_para_text = normalize_for_dedup(txt)
                        sentences_count = len(re.findall(r'[.!?。！？]', txt))
                        if sentences_count >= 2 or len(clean_para_text) > 60:
                            is_duplicate = False
                            for pool_text in knowledge_paragraphs_pool:
                                if clean_para_text in pool_text or pool_text in clean_para_text:
                                    is_duplicate = True
                                    break
                            if is_duplicate:
                                continue 
                        slide_lines.append(txt)
            
            skip_indices = set()
            for w_idx, line in enumerate(slide_lines):
                if w_idx in skip_indices:
                    continue
                
                if re.match(r'^\s*\d+\s*[a-zA-Z]', line):
                    p_w = doc.add_paragraph()
                    p_w.paragraph_format.line_spacing = 1.5
                    p_w.paragraph_format.space_before = Pt(0)
                    p_w.paragraph_format.space_after = Pt(0)
                    
                    r_vocab = p_w.add_run(line)
                    set_run_font(r_vocab, 10.5, bold=True) 
                    
                    search_idx = w_idx + 1
                    while search_idx < len(slide_lines):
                        next_line = slide_lines[search_idx]
                        if re.match(r'^\s*\d+\s*[a-zA-Z]', next_line):
                            break
                        if re.search(r'[\u4e00-\u9fa5]', next_line):
                            p_w.add_run("    ") 
                            r_trans = p_w.add_run(next_line)
                            set_run_font(r_trans, 10.5, bold=False)
                            r_trans.font.color.rgb = RGBColor(0, 0, 255) 
                            skip_indices.add(search_idx)
                            break
                        search_idx += 1
                else:
                    p_normal = doc.add_paragraph()
                    p_normal.paragraph_format.line_spacing = 1.5
                    p_normal.paragraph_format.space_before = Pt(0)
                    p_normal.paragraph_format.space_after = Pt(0)
                    
                    lines = line.split('\n')
                    for l_idx, l_text in enumerate(lines):
                        add_runs_with_blue_tags(p_normal, l_text)
                        if l_idx < len(lines) - 1:
                            p_normal.add_run('\n')
                            
            last_active_name = active_name
            continue

    doc.save(output_path)
    print(f"✅ 英语讲义处理完成: {output_path}\n")

if __name__ == '__main__':
    if not os.path.exists(ICON_DIR): os.makedirs(ICON_DIR)
    files = [f for f in os.listdir('.') if f.lower().endswith('.pptx') and not f.startswith('~$')]
    for f in files: process_single_english_ppt(f, DOC_TEMPLATE, ICON_DIR)
