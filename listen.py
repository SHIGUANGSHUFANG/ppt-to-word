import os
import re
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from docx import Document
from docx.shared import Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

# ==========================================
# 1. 基础净化与排版引擎
# ==========================================
def clean_filename(filename):
    return re.sub(r'[^\w\s\.\-\(\)]', '_', filename)

def set_run_font(run, size_pt, bold=False):
    run.font.name = 'Times New Roman'
    rPr = run._element.get_or_add_rPr()
    rFonts = rPr.get_or_add_rFonts()
    rFonts.set(qn('w:eastAsia'), '宋体')
    rFonts.set(qn('w:ascii'), 'Times New Roman')
    rFonts.set(qn('w:hAnsi'), 'Times New Roman')
    run.font.size = Pt(size_pt)
    run.bold = bold

def add_compact_paragraph(doc, text, bold=False):
    if text != "##FORCE_EMPTY##":
        text = text.replace('\xa0', ' ').strip()
        if not text: return None
        
    p = doc.add_paragraph()
    p.paragraph_format.line_spacing = 1.5
    p.paragraph_format.space_before = Pt(0)
    p.paragraph_format.space_after = Pt(0)
    
    if text and text != "##FORCE_EMPTY##":
        run = p.add_run(text)
        set_run_font(run, 10.5, bold=bold)
    return p

# ==========================================
# 表格与边框绘制
# ==========================================
def set_table_borders(table):
    tbl = table._element
    tblPr = tbl.xpath('w:tblPr')
    if not tblPr:
        tblPr = OxmlElement('w:tblPr')
        tbl.insert(0, tblPr)
    else:
        tblPr = tblPr[0]
        
    tblBorders = tblPr.xpath('w:tblBorders')
    if not tblBorders:
        tblBorders = OxmlElement('w:tblBorders')
        tblPr.append(tblBorders)
    else:
        tblBorders = tblBorders[0]
        tblBorders.clear()
        
    for b_name in ['top', 'left', 'bottom', 'right', 'insideH', 'insideV']:
        b = OxmlElement(f'w:{b_name}')
        b.set(qn('w:val'), 'single')
        b.set(qn('w:sz'), '4')  
        b.set(qn('w:space'), '0')
        b.set(qn('w:color'), 'auto')
        tblBorders.append(b)

# ==========================================
# 2. 空间雷达与核心提取器
# ==========================================
def flatten_shapes(shapes):
    res = []
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            res.extend(flatten_shapes(shape.shapes))
        else:
            res.append(shape)
    return res

def get_shape_text(shape):
    if getattr(shape, 'has_text_frame', False):
        return shape.text
    elif getattr(shape, 'has_table', False):
        return "".join(cell.text for row in shape.table.rows for cell in row.cells)
    return ""

def get_ignored_shape_ids(shapes):
    ignore_ids = set()
    valid_shapes = [x for x in shapes if getattr(x, 'has_text_frame', False) or getattr(x, 'has_table', False)]
    for shape in valid_shapes:
        text_clean = re.sub(r'\s+', '', get_shape_text(shape))
        if "原文重现" in text_clean:
            ignore_ids.add(id(shape))
            for rs in get_shapes_to_right(shape, valid_shapes):
                ignore_ids.add(id(rs))
    return ignore_ids

def get_shapes_to_right(shape_a, all_shapes):
    right_shapes = []
    if not hasattr(shape_a, 'top') or not hasattr(shape_a, 'left'): return []
    for b in all_shapes:
        if b == shape_a: continue
        if not hasattr(b, 'top') or not hasattr(b, 'left'): continue
        if not (shape_a.top + shape_a.height < b.top or shape_a.top > b.top + b.height):
            if b.left > shape_a.left:
                right_shapes.append(b)
    return right_shapes

def stitch_and_bold_vocab(raw_text, pending_number=None, sec_title=""):
    raw_text = raw_text.replace('\x0b', '\n').replace('\r', '\n').replace('\xa0', ' ')
    raw_lines = [line.strip() for line in raw_text.split('\n') if line.strip()]
    
    headers = []
    formatted = []
    
    for line in raw_lines:
        clean_l = re.sub(r'\s+', '', line)
        if clean_l in ["词汇精讲", "用法搭配", "原文重现", "重点词汇", "听力原文", "听力文本"]:
            continue
        if sec_title and re.sub(r'\s+', '', sec_title) in clean_l:
            continue
            
        if pending_number:
            line = f"{pending_number}  {line}"
            pending_number = None
            
        if re.match(r'^\d{1,2}[\.、]?$', line):
            pending_number = line
        else:
            is_bold = bool(re.match(r'^\d{1,2}[\.、]?\s*', line) and re.search(r'[\u4e00-\u9fa5]', line))
            if is_bold:
                headers.append(line)
            formatted.append((line, is_bold))
            
    return headers, formatted, pending_number

def extract_table_vocab(doc, ppt_table, font_size=10.5, pending_number=None, sec_title=""):
    word_table = doc.add_table(rows=len(ppt_table.rows), cols=len(ppt_table.columns))
    set_table_borders(word_table)
    headers_list = []
    
    for r_idx, row in enumerate(ppt_table.rows):
        for c_idx, cell in enumerate(row.cells):
            is_duplicate = False
            if r_idx > 0 and cell.text == ppt_table.cell(r_idx-1, c_idx).text and cell.text.strip():
                is_duplicate = True
            if c_idx > 0 and cell.text == ppt_table.cell(r_idx, c_idx-1).text and cell.text.strip():
                is_duplicate = True
            
            word_cell = word_table.cell(r_idx, c_idx)
            p = word_cell.paragraphs[0]
            p.paragraph_format.line_spacing = 1.5
            
            if not is_duplicate:
                headers, formatted, pending_number = stitch_and_bold_vocab(cell.text, pending_number, sec_title)
                headers_list.extend(headers)
                
                for i, (text, bold) in enumerate(formatted):
                    if i > 0:
                        p = word_cell.add_paragraph()
                        p.paragraph_format.line_spacing = 1.5
                    run = p.add_run(text)
                    set_run_font(run, font_size, bold=bold)
                    
    return headers_list, pending_number

def extract_table_normal(doc, ppt_table, font_size=10.5, sec_title=""):
    word_table = doc.add_table(rows=len(ppt_table.rows), cols=len(ppt_table.columns))
    set_table_borders(word_table)
    for r_idx, row in enumerate(ppt_table.rows):
        for c_idx, cell in enumerate(row.cells):
            is_duplicate = False
            if r_idx > 0 and cell.text == ppt_table.cell(r_idx-1, c_idx).text and cell.text.strip():
                is_duplicate = True
            if c_idx > 0 and cell.text == ppt_table.cell(r_idx, c_idx-1).text and cell.text.strip():
                is_duplicate = True
            
            word_cell = word_table.cell(r_idx, c_idx)
            p = word_cell.paragraphs[0]
            p.paragraph_format.line_spacing = 1.5
            
            if not is_duplicate:
                raw_lines = cell.text.split('\n')
                cleaned_lines = []
                for line in raw_lines:
                    l_strip = line.replace('\xa0', ' ').strip()
                    if l_strip and l_strip not in ["知识精讲", sec_title]:
                        cleaned_lines.append(l_strip)
                
                for i, text in enumerate(cleaned_lines):
                    if i > 0:
                        p = word_cell.add_paragraph()
                        p.paragraph_format.line_spacing = 1.5
                    run = p.add_run(text)
                    set_run_font(run, font_size)

def get_slide_sections_robust(prs):
    slide_sections = {}
    target_keys = ["短对话", "长对话", "独白"]
    try:
        sections = prs.element.xpath('.//*[local-name()="section"]')
        slide_id_to_section = {}
        for section in sections:
            name = section.get('name')
            for sld in section.xpath('.//*[local-name()="sldId"]'):
                val = sld.get('id')
                if val: slide_id_to_section[int(val)] = name
        for idx, slide in enumerate(prs.slides, 1):
            slide_sections[idx] = slide_id_to_section.get(slide.slide_id, "")
    except: pass
    
    if not any(any(k in v for v in slide_sections.values()) for k in target_keys):
        current_section = ""
        for idx, slide in enumerate(prs.slides, 1):
            found_key = None
            title_txt = slide.shapes.title.text if slide.shapes.title and slide.shapes.title.has_text_frame else ""
            l_name = slide.slide_layout.name
            for key in target_keys:
                if key in title_txt or key in l_name: found_key = key; break
            if not found_key:
                for shape in slide.shapes:
                    if getattr(shape, 'has_text_frame', False):
                        txt = re.sub(r'\s+', '', shape.text)
                        for key in target_keys:
                            if key in txt and len(txt) < 15: found_key = key; break
                    if found_key: break
            if found_key: current_section = found_key
            slide_sections[idx] = current_section
    return slide_sections

def extract_title_bold_only(shape):
    if not getattr(shape, 'has_text_frame', False): return ""
    bold_pieces = []
    for paragraph in shape.text_frame.paragraphs:
        p_txt = ""
        for run in paragraph.runs:
            if run.font.bold: p_txt += run.text
        if p_txt.strip(): bold_pieces.append(p_txt.strip())
    return " ".join(bold_pieces).strip()

def check_target_text_in_slide(slide, target_text):
    for shape in flatten_shapes(slide.shapes):
        if target_text in re.sub(r'\s+', '', get_shape_text(shape)): return True
    return False

def check_yuanwen_in_slide(slide):
    for shape in flatten_shapes(slide.shapes):
        txt_clean = re.sub(r'\s+', '', get_shape_text(shape)).replace("原文重现", "")
        if "听力文本" in txt_clean or "原文" in txt_clean: return True
    return False

# ==========================================
# 3. 英语听力核心调度引擎 (封装版)
# ==========================================
# 👇 注意这里的函数名和参数，已经和 app.py 中的调用完全对齐
def process_single_listening_ppt(ppt_path, doc_template, icon_dir, output_path):
    prs = Presentation(ppt_path)
    doc = Document(doc_template)
    
    slide_sections = get_slide_sections_robust(prs)
    target_sections = {"短对话": "一、短对话", "长对话": "二、长对话", "独白": "三、独白"}

    # --- 模块 0：【全局透视】提前收集所有重点词汇 (全地形穿透) ---
    global_vocab_headers = []
    for i, slide in enumerate(prs.slides, 1):
        sec = slide_sections.get(i, "")
        sec_title = target_sections.get(sec, "")
        
        if check_target_text_in_slide(slide, "用法搭配"):
            shapes = flatten_shapes(slide.shapes)
            ignore_ids = get_ignored_shape_ids(shapes)
            valid_shapes = [s for s in shapes if id(s) not in ignore_ids and (getattr(s, 'has_text_frame', False) or getattr(s, 'has_table', False))]
            valid_shapes.sort(key=lambda x: (getattr(x, 'top', 0) or 0, getattr(x, 'left', 0) or 0))
            
            pending_number = None
            for shape in valid_shapes:
                if getattr(shape, 'has_table', False):
                    for row in shape.table.rows:
                        for cell in row.cells:
                            headers, _, pending_number = stitch_and_bold_vocab(cell.text, pending_number, sec_title)
                            global_vocab_headers.extend(headers)
                else:
                    headers, _, pending_number = stitch_and_bold_vocab(shape.text, pending_number, sec_title)
                    global_vocab_headers.extend(headers)

    # --- 模块 1：标题处理 ---
    for slide in prs.slides:
        if slide.slide_layout.name.strip() == "标题":
            title_text = ""
            for shape in slide.shapes:
                t = extract_title_bold_only(shape)
                if t: title_text = t; break
            
            if title_text:
                merged_title = f"##INSERT_IMAGE:标题块.emf##{title_text}"
                if len(doc.paragraphs) > 0 and not doc.paragraphs[0].text.strip():
                    p = doc.paragraphs[0]
                    p.text = merged_title
                else:
                    p = doc.paragraphs[0].insert_paragraph_before(merged_title) if len(doc.paragraphs) > 0 else doc.add_paragraph(merged_title)
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                for run in p.runs: set_run_font(run, 20, bold=True)
            break 

    # --- 模块 2：课堂目标处理 ---
    for slide in prs.slides:
        if slide.slide_layout.name.strip() == "知识目标":
            add_compact_paragraph(doc, "##INSERT_IMAGE:课堂目标.emf##")
            counter = 1
            shapes = flatten_shapes(slide.shapes)
            shapes.sort(key=lambda s: getattr(s, 'top', 0) or 0)
            
            for shape in shapes:
                if not getattr(shape, 'has_text_frame', False): continue
                if shape == slide.shapes.title: continue 
                
                for paragraph in shape.text_frame.paragraphs:
                    if not paragraph.text.strip(): continue
                    p = doc.add_paragraph()
                    p.paragraph_format.line_spacing = 1.5
                    
                    r_num = p.add_run(f"{counter}. ")
                    set_run_font(r_num, 10.5, bold=False)
                    
                    for run in paragraph.runs:
                        text = run.text
                        if not text.strip(): 
                            p.add_run(text); continue
                            
                        is_black = True
                        if run.font and run.font.color and hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                            rgb_str = str(run.font.color.rgb)
                            if rgb_str != "000000" and rgb_str.lower() != "none":
                                is_black = False
                        
                        r_word = p.add_run(text)
                        set_run_font(r_word, 10.5, bold=not is_black)
                    counter += 1
            break

    # --- 模块 2.5：【课前预习板块】 ---
    if global_vocab_headers:
        add_compact_paragraph(doc, "##INSERT_IMAGE:课前预习.emf##")
        add_compact_paragraph(doc, "听伴学本上的音频，并在讲义横线上填写中文含义。", bold=True)
        
        for i, header in enumerate(global_vocab_headers):
            add_compact_paragraph(doc, header, bold=True)
            if i < len(global_vocab_headers) - 1:
                add_compact_paragraph(doc, "##FORCE_EMPTY##")

    # --- 模块 3：按节提取听力内容 ---
    appendix_buffer = {"短对话": [], "长对话": [], "独白": []}
    
    sections_order = []
    seen = set()
    for i in range(1, len(prs.slides) + 1):
        sec = slide_sections.get(i, "")
        for k in target_sections.keys():
            if k in sec: 
                sec = k
                break
        if sec in target_sections and sec not in seen:
            sections_order.append(sec)
            seen.add(sec)

    for sec in sections_order:
        sec_slides = [i for i in range(1, len(prs.slides)+1) if sec in slide_sections.get(i, "")]
        zhishi_icon_inserted = False
        last_main_text_type = None
        sec_title = target_sections.get(sec, "")
        
        for i in sec_slides:
            s = prs.slides[i-1]
            
            is_yuanwen = check_yuanwen_in_slide(s)
            is_cihui = check_target_text_in_slide(s, "用法搭配")
            
            if is_yuanwen:
                shapes = flatten_shapes(s.shapes)
                shapes.sort(key=lambda x: getattr(x, 'top', 0) or 0)
                
                has_long_yuanwen_shape = False
                for shape in shapes:
                    txt = get_shape_text(shape).strip()
                    txt_clean = re.sub(r'\s+', '', txt).replace("原文重现", "")
                    if ("听力文本" in txt_clean or "原文" in txt_clean) and len(txt) > 15:
                        has_long_yuanwen_shape = True
                        break
                
                for shape in shapes:
                    txt = get_shape_text(shape).strip()
                    if not txt: continue
                    txt_clean = re.sub(r'\s+', '', txt).replace("原文重现", "")
                    
                    if has_long_yuanwen_shape:
                        if "听力文本" in txt_clean or "原文" in txt_clean:
                            appendix_buffer[sec].append(txt)
                    else:
                        if "听力文本" in txt_clean or "原文" in txt_clean or len(txt) > 15:
                            if txt not in [sec, sec_title]:
                                appendix_buffer[sec].append(txt)

            elif is_cihui:
                if last_main_text_type == "cihui":
                    p_sep = doc.add_paragraph()
                    p_sep.paragraph_format.line_spacing = 1.5
                    run_sep = p_sep.add_run("﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌")
                    set_run_font(run_sep, 10.5)
                    p_sep.alignment = WD_ALIGN_PARAGRAPH.CENTER
                
                last_main_text_type = "cihui"
                add_compact_paragraph(doc, "##INSERT_IMAGE:词汇精讲.emf##")
                
                shapes = flatten_shapes(s.shapes)
                ignore_ids = get_ignored_shape_ids(shapes)
                valid_shapes = [x for x in shapes if id(x) not in ignore_ids and (getattr(x, 'has_text_frame', False) or getattr(x, 'has_table', False))]
                valid_shapes.sort(key=lambda x: (getattr(x, 'top', 0) or 0, getattr(x, 'left', 0) or 0))

                pending_number = None
                for shape in valid_shapes:
                    if getattr(shape, 'has_table', False):
                        _, pending_number = extract_table_vocab(doc, shape.table, pending_number=pending_number, sec_title=sec_title)
                    else:
                        headers, formatted, pending_number = stitch_and_bold_vocab(shape.text, pending_number, sec_title)
                        for line, is_bold in formatted:
                            add_compact_paragraph(doc, line, bold=is_bold)
                            
                if pending_number:
                    add_compact_paragraph(doc, pending_number)

            else:
                last_main_text_type = "zhishi"
                if not zhishi_icon_inserted:
                    add_compact_paragraph(doc, "##INSERT_IMAGE:知识精讲.emf##")
                    add_compact_paragraph(doc, sec_title, bold=True)
                    zhishi_icon_inserted = True
                
                shapes = flatten_shapes(s.shapes)
                shapes.sort(key=lambda x: getattr(x, 'top', 0) or 0)
                
                for shape in shapes:
                    if getattr(shape, 'has_table', False):
                        extract_table_normal(doc, shape.table, sec_title=sec_title)
                    elif getattr(shape, 'has_text_frame', False):
                        raw_lines = [l.strip() for l in shape.text.replace('\x0b','\n').replace('\r','\n').split('\n') if l.strip()]
                        for line in raw_lines:
                            if line and line not in ["知识精讲", sec, sec_title]:
                                add_compact_paragraph(doc, line)

    # ==========================================
    # 4. 附录生成器
    # ==========================================
    has_appendix = any(len(appendix_buffer[k]) > 0 for k in appendix_buffer)
    if has_appendix:
        p_sep = doc.add_paragraph()
        p_sep.paragraph_format.line_spacing = 1.5
        p_sep.paragraph_format.space_before = Pt(12)
        run_sep = p_sep.add_run("﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌﹌")
        set_run_font(run_sep, 10.5)
        p_sep.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        p_app_title = add_compact_paragraph(doc, "【附录】听力原文", bold=True)
        if p_app_title: p_app_title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        for sec in sections_order:
            if appendix_buffer[sec]:
                add_compact_paragraph(doc, target_sections[sec], bold=True)
                counter = 1
                for text_block in appendix_buffer[sec]:
                    lines = text_block.split('\n')
                    first_line = True
                    printed_anything = False
                    for line in lines:
                        line = line.strip()
                        if not line: continue
                        
                        if first_line:
                            line = re.sub(r'^(【听力原文】|【原文】|【听力文本】|听力原文|听力文本|原文)[:：]?\s*', '', line)
                            if line:
                                add_compact_paragraph(doc, f"{counter}. {line}")
                                first_line = False
                                printed_anything = True
                        else:
                            add_compact_paragraph(doc, line)
                            printed_anything = True
                            
                    if printed_anything:
                        counter += 1

    # ==========================================
    # 5. 全局清理引擎
    # ==========================================
    while len(doc.paragraphs) > 0 and not doc.paragraphs[0].text.strip() and len(doc.paragraphs[0].runs) == 0:
        p_first = doc.paragraphs[0]._element
        p_first.getparent().remove(p_first)

    consecutive_empty = 0
    p_to_remove = []
    for p in doc.paragraphs:
        if not p.text.strip() and (len(p.runs) == 0 or all(not r.text.strip() for r in p.runs)):
            consecutive_empty += 1
            if consecutive_empty > 2:  
                p_to_remove.append(p)
        else:
            consecutive_empty = 0
            
    for p in p_to_remove:
        pPr = p._element
        pPr.getparent().remove(pPr)

    # 最终保存到 app.py 传递过来的路径
    doc.save(output_path)