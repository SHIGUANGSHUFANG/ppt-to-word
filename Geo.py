import csv
import re
import os
import io
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from docx import Document
from docx.shared import Emu, Pt
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

def is_option_text(text):
    text = text.strip()
    if re.match(r'^[A-D][\.、．\s]', text):
        return True
    if "A." in text and "B." in text:
        return True
    return False

def add_paragraph_with_style(doc, text, style_name):
    text = text.strip()
    if not text: return
    if style_name:
        try:
            doc.add_paragraph(text, style=style_name)
        except KeyError:
            doc.add_paragraph(text)
    else:
        doc.add_paragraph(text)

def set_table_borders(table):
    tbl = table._element
    tblPr = tbl.xpath('w:tblPr')
    if not tblPr:
        tblPr = OxmlElement('w:tblPr')
        tbl.insert(0, tblPr)
    else:
        tblPr = tblPr[0]
        
    tblBorders = tblPr.xpath('w:tblBorders')
    if not tblBorders:
        tblBorders = OxmlElement('w:tblBorders')
        tblPr.append(tblBorders)
    else:
        tblBorders = tblBorders[0]
        tblBorders.clear()
        
    for b_name in ['top', 'left', 'bottom', 'right', 'insideH', 'insideV']:
        b = OxmlElement(f'w:{b_name}')
        b.set(qn('w:val'), 'single')
        b.set(qn('w:sz'), '4')  
        b.set(qn('w:space'), '0')
        b.set(qn('w:color'), 'auto')
        tblBorders.append(b)

# 【核心重构】：把 main() 改成了带 4 个参数的标准接口
def process_single_physics_ppt(ppt_path, doc_template, icon_dir, output_path=None):
    csv_path = '模板-样式对应关系.csv'

    if not os.path.exists(icon_dir):
        os.makedirs(icon_dir)

    # 兼容网页端传来的路径
    if output_path is None:
        raw_filename = os.path.splitext(os.path.basename(ppt_path))[0]
        output_path = f"{raw_filename}物理讲义.docx"

    config = {}
    try:
        with open(csv_path, mode='r', encoding='utf-8-sig') as f:
            reader = csv.DictReader(f)
            for row in reader:
                config[row['幻灯片版式'].strip()] = row
    except Exception as e:
        raise Exception(f"❌ 读取 CSV 失败，请确认 Github 仓库中包含 {csv_path} 文件！报错信息: {e}")

    prs = Presentation(ppt_path)
    doc = Document(doc_template)
    last_layout_name = None
    MAX_WORD_WIDTH_EMU = 5257800 

    for slide_number, slide in enumerate(prs.slides, 1):
        layout_name = slide.slide_layout.name.strip()
        cfg = config.get(layout_name)

        if not cfg:
            continue
            
        print(f"正在处理第 {slide_number} 页：{layout_name}")

        if layout_name != last_layout_name:
            icon_file = cfg.get('插入矢量图名称', '').strip()
            if icon_file and icon_file != '不插入':
                p = doc.add_paragraph()
                pic_style = cfg.get('图片部分word样式名称')
                if pic_style:
                    try: p.style = pic_style
                    except: pass
                p.add_run(f"##INSERT_IMAGE:{icon_file}##")

        extract_rule = cfg.get('提取内容', '')
        shapes = [s for s in slide.shapes if hasattr(s, 'top') and s.top is not None]
        shapes.sort(key=lambda s: s.top)

        for shape in shapes:
            is_title_shape = (shape == slide.shapes.title)
            
            if getattr(shape, 'has_table', False):
                if '全部文本' in extract_rule or '仅文本' in extract_rule:
                    ppt_table = shape.table
                    word_table = doc.add_table(rows=len(ppt_table.rows), cols=len(ppt_table.columns))
                    
                    set_table_borders(word_table)

                    for r_idx, row in enumerate(ppt_table.rows):
                        for c_idx, cell in enumerate(row.cells):
                            is_duplicate = False
                            if r_idx > 0 and cell.text == ppt_table.cell(r_idx-1, c_idx).text and cell.text.strip():
                                is_duplicate = True
                            if c_idx > 0 and cell.text == ppt_table.cell(r_idx, c_idx-1).text and cell.text.strip():
                                is_duplicate = True

                            word_cell = word_table.cell(r_idx, c_idx)
                            p = word_cell.paragraphs[0]
                            p.text = ""
                            
                            target_style = cfg.get('非标题部分word样式名称')
                            if target_style:
                                try: p.style = target_style
                                except: pass

                            if not is_duplicate:
                                raw_lines = cell.text.split('\n')
                                cleaned_lines = [line.replace('\xa0', ' ').strip() for line in raw_lines if line.strip()]
                                if cleaned_lines:
                                    p.add_run('\n'.join(cleaned_lines))

            elif shape.has_text_frame:
                if '全部文本' in extract_rule or '仅文本' in extract_rule or (is_title_shape and '标题' in extract_rule):
                    raw_lines = shape.text.split('\n')
                    cleaned_lines = [line.replace('\xa0', ' ').strip() for line in raw_lines if line.strip()]
                    if cleaned_lines:
                        text_content = '\n'.join(cleaned_lines)
                        target_style = ""
                        if is_title_shape: target_style = cfg.get('标题部分word样式名称')
                        elif is_option_text(text_content): target_style = cfg.get('试题选项部分word样式名称')
                        elif "例题" in layout_name or "练习" in layout_name: target_style = cfg.get('试题题干部分word样式名称')
                        else: target_style = cfg.get('非标题部分word样式名称')
                        add_paragraph_with_style(doc, text_content, target_style)

            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE and '图片' in extract_rule:
                try:
                    image = shape.image
                    ext, image_bytes = image.ext.lower(), image.blob
                    sw, sh = shape.width, shape.height
                    if sw > MAX_WORD_WIDTH_EMU:
                        ratio = MAX_WORD_WIDTH_EMU / sw
                        sw, sh = MAX_WORD_WIDTH_EMU, int(sh * ratio)

                    p = doc.add_paragraph()
                    pic_style = cfg.get('图片部分word样式名称')
                    if pic_style:
                        try: p.style = pic_style
                        except: pass

                    if ext in ['emf', 'wmf']:
                        img_filename = f"ppt_extract_s{slide_number}_{shape.shape_id}.{ext}"
                        with open(os.path.join(icon_dir, img_filename), 'wb') as f: f.write(image_bytes)
                        p.add_run(f"##INSERT_IMAGE:{img_filename}|{sw/12700:.1f}|{sh/12700:.1f}##")
                    else:
                        r = p.add_run()
                        r.add_picture(io.BytesIO(image_bytes), width=Emu(sw), height=Emu(sh))
                except Exception as e: print(f"   ⚠️ 提取插图失败: {e}")

        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide.notes_text_frame is not None:
                notes_text = notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    raw_notes_lines = notes_text.split('\n')
                    cleaned_notes = [line.replace('\xa0', ' ').strip() for line in raw_notes_lines if line.strip()]
                    
                    if cleaned_notes:
                        final_notes = '\n'.join(cleaned_notes)
                        notes_style = cfg.get('备注部分word样式名称') 
                        if not notes_style:
                            notes_style = cfg.get('非标题部分word样式名称')
                        
                        add_paragraph_with_style(doc, final_notes, notes_style)

        last_layout_name = layout_name

    doc.save(output_path)
    print(f"\n✅ 物理讲义处理完成！已生成文件：{output_path}")

if __name__ == '__main__':
    files = [f for f in os.listdir('.') if f.lower().endswith('.pptx') and not f.startswith('~$')]
    for f in files: process_single_physics_ppt(f, '模板.docx', 'icons')